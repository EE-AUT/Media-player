import re
from docx import Document
import os
from pptx import Presentation
import glob
import csv


# read data from csv file
def read_csv(filename):
    try:
        tagsDict = {}
        movie_Session = ""
        output = {}
        with open(filename) as csvfile:
            file_reader = csv.reader(csvfile, delimiter='#')
            for row in file_reader:
                # classification data with len and find tag session and tags
                if len(row) == 1:
                    output.update({movie_Session: tagsDict})
                    tagsDict = {}
                    movie_Session = row[0]
                elif len(row) == 2:
                    tagsDict.update({row[0]: row[1]})
                elif len(row) == 3:
                    tagsDict.update({row[0]: row[1]})  # read from bookmarks
            output.update({movie_Session: tagsDict})

            del output['']  # delete empty key

            return output
            # handle Exception
    except:
        return {}


# read data from docx file
def read_docx(filename):
    try:
        # read from word tables
        wordDoc = Document(filename)
        output = {}
        tagsDict = {}
        name = []
        movie_Session = ""
        for table in wordDoc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text != "":  # if cell has text
                        # classification data and find session and tags
                        if re.search("Session = ", cell.text):
                            try:
                                # create dictionary
                                tagsDict = {name[i]: name[i+1]
                                            for i in range(0, len(name), 2)}
                            except:
                                tagsDict = {}
                            output.update({movie_Session: tagsDict})
                            name = []
                            movie_Session = re.split(
                                "Session = ", cell.text)[1]

                        else:
                            name.append(cell.text)

        # convert list to dict
        tagsDict = {name[i]: name[i+1] for i in range(0, len(name), 2)}
        output.update({movie_Session: tagsDict})

        del output[""]  # delete empty key
        return output
    except:  # handle exception
        return {}


# get data from PPtx files
def read_pptx(filename):
    try:
        output = {}
        # open presentation
        for eachfile in glob.glob(filename):
            prs = Presentation(eachfile)
            movie_Session = ""
            movie_string = ""

            a = {}
            # iterate each shape to find text
            for slide in prs.slides:  # get slides
                for shape in slide.shapes:
                    if hasattr(shape, "text"):  # if there is shape in tags
                        if re.search("Session = #", shape.text):
                            # update dictionary
                            a.update({movie_Session: movie_string})
                            movie_string = ""
                            # find session name with our role
                            movie_Session = re.split(
                                "Session = #", shape.text)[-1]
                        else:
                            movie_string = movie_string + "\n" + shape.text
                a.update({movie_Session: movie_string})

        for key in a:
            output.update({key: _split_text(a[key])})
        del output['']
        return output
    except:
        return {}


def _split_text(text):
    split_Enter = re.split(r"\n", text)
    split_Enter = [item for item in split_Enter if item != ""]

    output = {}

    for item in split_Enter:  # split text by Enter(\n)
        if re.search("\d\d:\d\d:\d\d", item):  # first format for time
            time = re.search("\d\d:\d\d:\d\d", item).group()
            try:
                text = re.split("\d\d:\d\d:\d\d:", item)[1]
                output.update({text: time})
            except:
                pass

        elif re.search("\d\d:\d\d", item):  # second format for time
            time = re.search("\d\d:\d\d", item).group()
            try:
                text = re.split("\d\d:\d\d:", item)[1]
                output.update({text: time})
            except:
                pass

    return output


# edit tag using re module
def edit_Tags(befor, after, filename):
    with open(filename) as csvfile:
        string = csvfile.read()

    with open(filename, "w") as csvfile:
        csvfile.write(re.sub(befor, after, string))

